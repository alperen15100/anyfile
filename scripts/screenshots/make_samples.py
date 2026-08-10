#!/usr/bin/env python3
"""
Generates realistic but entirely fictional sample documents for Gander's
Play Store screenshots. Every company, person, address and figure here is
invented; nothing imitates a real organisation.

Outputs into ./out :
  Tenancy Agreement - 14 Alder Court.pdf   (6 page residential lease)
  Q3 Operating Budget FY26.xlsx            (5 sheet workbook)
  Field Survey Report - Willowmere.docx    (formatted report)
  Willowmere Kickoff.pptx                  (deck)
"""
import os
from pathlib import Path

OUT = Path(__file__).parent / "out"
OUT.mkdir(exist_ok=True)

# --------------------------------------------------------------------------
# PDF: residential tenancy agreement
# --------------------------------------------------------------------------
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import mm
from reportlab.lib import colors
from reportlab.platypus import (BaseDocTemplate, Frame, PageTemplate, Paragraph,
                                Spacer, Table, TableStyle, KeepTogether)

INK = colors.HexColor("#1a1a1a")
RULE = colors.HexColor("#c8c8c8")
ACCENT = colors.HexColor("#1F3A5F")

ss = getSampleStyleSheet()
S = {}
S["title"] = ParagraphStyle("t", parent=ss["Normal"], fontName="Times-Bold",
                            fontSize=19, leading=23, alignment=1,
                            textColor=INK, spaceAfter=3)
S["sub"] = ParagraphStyle("s", parent=ss["Normal"], fontName="Times-Roman",
                          fontSize=10.5, leading=14, alignment=1,
                          textColor=colors.HexColor("#4a4a4a"), spaceAfter=14)
S["h"] = ParagraphStyle("h", parent=ss["Normal"], fontName="Times-Bold",
                        fontSize=11.5, leading=14, textColor=ACCENT,
                        spaceBefore=11, spaceAfter=5)
S["body"] = ParagraphStyle("b", parent=ss["Normal"], fontName="Times-Roman",
                           fontSize=10, leading=14.5, alignment=4,
                           textColor=INK, spaceAfter=7)
S["sub_c"] = ParagraphStyle("sc", parent=S["body"], leftIndent=16, spaceAfter=5)
S["cell"] = ParagraphStyle("c", parent=ss["Normal"], fontName="Times-Roman",
                           fontSize=9.5, leading=13, textColor=INK)
S["cellb"] = ParagraphStyle("cb", parent=S["cell"], fontName="Times-Bold")


def lease_pdf():
    """Superseded by make_lease.py, which writes the same filename.

    This reportlab version came first. It lays the parties block out as a real
    table, which trips the WebView rendering fault described in the README, so
    the shipped lease is the python-docx one. Run make_lease.py after this and it
    overwrites the output. Kept because it is the smaller reproduction of that
    fault if anyone wants to chase it.
    """
    path = OUT / "Tenancy Agreement - 14 Alder Court.pdf"

    def decor(canvas, doc):
        canvas.saveState()
        canvas.setFont("Times-Roman", 8)
        canvas.setFillColor(colors.HexColor("#777777"))
        canvas.drawString(20 * mm, 13 * mm,
                          "Assured Shorthold Tenancy — 14 Alder Court, Kestrel Bay KB7 2QN")
        canvas.drawRightString(A4[0] - 20 * mm, 13 * mm, "Page %d of 6" % doc.page)
        canvas.setStrokeColor(RULE)
        canvas.setLineWidth(0.5)
        canvas.line(20 * mm, 16.5 * mm, A4[0] - 20 * mm, 16.5 * mm)
        canvas.drawRightString(A4[0] - 20 * mm, A4[1] - 14 * mm, "REF: AC-14/2026-08")
        canvas.restoreState()

    doc = BaseDocTemplate(str(path), pagesize=A4,
                          leftMargin=20 * mm, rightMargin=20 * mm,
                          topMargin=19 * mm, bottomMargin=20 * mm,
                          title="Tenancy Agreement - 14 Alder Court",
                          author="Kestrel Bay Lettings")
    frame = Frame(doc.leftMargin, doc.bottomMargin, doc.width, doc.height, id="f")
    doc.addPageTemplates([PageTemplate(id="p", frames=[frame], onPage=decor)])

    f = []
    f.append(Paragraph("RESIDENTIAL TENANCY AGREEMENT", S["title"]))
    f.append(Paragraph("Assured Shorthold Tenancy under Part I, Chapter II of the "
                       "Housing Act (fictional jurisdiction of Kestrel Bay)", S["sub"]))

    parties = [
        [Paragraph("<b>Landlord</b>", S["cellb"]),
         # No registration number here, deliberately. See "Read this before
         # inventing an identifier" in the README.
         Paragraph("Marisol Okonkwo-Reyes, of 3 Thornhill Row, Kestrel Bay KB4 8LT, "
                   "acting through her managing agent, Aldergate Property Care.",
                   S["cell"])],
        [Paragraph("<b>Tenants</b>", S["cellb"]),
         Paragraph("Dev Raghunathan and Priya Raghunathan, jointly and severally "
                   "liable for the whole of the rent and for every obligation in "
                   "this agreement.", S["cell"])],
        [Paragraph("<b>Property</b>", S["cellb"]),
         Paragraph("The whole of the two-bedroom first-floor flat known as 14 Alder "
                   "Court, Marrow Lane, Kestrel Bay KB7 2QN, together with the "
                   "fixtures listed in the Inventory at Schedule 2 and the allocated "
                   "parking bay numbered 14.", S["cell"])],
        [Paragraph("<b>Term</b>", S["cellb"]),
         Paragraph("Twelve (12) months certain, beginning on 1 September 2026 and "
                   "ending on 31 August 2027, after which the tenancy continues as a "
                   "statutory periodic tenancy from month to month until ended under "
                   "clause 11.", S["cell"])],
        [Paragraph("<b>Rent</b>", S["cellb"]),
         Paragraph("£1,340 per calendar month, payable in advance on the 1st day "
                   "of each month by standing order to the client account named in "
                   "Schedule 1. The first payment falls due on 1 September 2026.",
                   S["cell"])],
        [Paragraph("<b>Deposit</b>", S["cellb"]),
         Paragraph("£1,546, being five weeks' rent, protected within 30 days of "
                   "receipt in an approved custodial scheme. The scheme reference is "
                   "issued with the prescribed information at Schedule 3.", S["cell"])],
    ]
    t = Table(parties, colWidths=[30 * mm, 140 * mm])
    t.setStyle(TableStyle([
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ("TOPPADDING", (0, 0), (-1, -1), 5),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
        ("LEFTPADDING", (0, 0), (-1, -1), 5),
        ("RIGHTPADDING", (0, 0), (-1, -1), 5),
        ("BACKGROUND", (0, 0), (0, -1), colors.HexColor("#EEF2F7")),
        ("GRID", (0, 0), (-1, -1), 0.5, RULE),
    ]))
    f.append(t)
    f.append(Spacer(1, 9))
    f.append(Paragraph(
        "It is agreed that the Landlord lets and the Tenants take the Property for the "
        "Term at the Rent, on the terms set out below. Where the Tenants are more than "
        "one person, every obligation binds each of them individually and all of them "
        "together. Words in the singular include the plural, and a reference to a "
        "clause is a reference to a clause of this agreement.", S["body"]))

    def clause(n, head, paras, subs=None):
        blk = [Paragraph("%d. %s" % (n, head), S["h"])]
        for p in paras:
            blk.append(Paragraph(p, S["body"]))
        if subs:
            for i, sp in enumerate(subs):
                blk.append(Paragraph("%d.%d&nbsp;&nbsp;%s" % (n, i + 1, sp), S["sub_c"]))
        return blk

    f += clause(1, "Rent and outgoings", [
        "The Tenants shall pay the Rent without deduction or set-off. If any part of "
        "the Rent remains unpaid for fourteen days after it became due, interest "
        "accrues on the unpaid amount at three per cent above the Bank of Kestrel "
        "base rate, calculated daily from the date the payment fell due until it is "
        "paid in full, whether or not the Landlord has formally demanded it.",
        "The Tenants shall pay all charges for electricity, gas, water, wastewater, "
        "broadband and council tax relating to the Property during the Term, "
        "including any standing charges and any reconnection fee arising from a "
        "supply cut off because of non-payment, and shall provide meter readings on "
        "the first and last day of the Term."])

    f += clause(2, "Condition and repair", [
        "The Tenants shall keep the interior of the Property in the same clean and "
        "tenantable condition in which they received it, fair wear and tear excepted, "
        "and shall make good any damage caused by them, by anyone living with them or "
        "by any visitor. The Landlord remains responsible for the structure and "
        "exterior, for the installations for the supply of water, gas, electricity "
        "and sanitation, and for the space and water heating."],
        subs=[
            "The Tenants shall report any disrepair, water ingress, or failure of a "
            "heating or sanitary installation to the managing agent within five "
            "working days of becoming aware of it. A defect that worsens because it "
            "was not reported may be charged to the Tenants to the extent of the "
            "additional cost.",
            "The Tenants shall test the smoke alarms monthly and shall not remove a "
            "battery from any alarm except to replace it immediately.",
            "The Tenants shall keep the Property adequately ventilated and heated so "
            "as to avoid condensation, and shall not dry laundry indoors without "
            "opening a window or running the extractor fan.",
            "The Tenants shall not redecorate, drill, or fix anything to a wall "
            "without the Landlord's written consent, which will not be unreasonably "
            "withheld for ordinary picture hooks.",
        ])

    f += clause(3, "Use of the Property", [
        "The Tenants shall use the Property as a single private dwelling for "
        "themselves and for no more than two other permitted occupiers named in "
        "Schedule 1. The Tenants shall not carry on any trade or business at the "
        "Property, other than quiet office work that generates no visitors, no "
        "signage and no deliveries beyond ordinary domestic volume.",
        "The Tenants shall not do anything at the Property that is a nuisance or "
        "annoyance to neighbours, that invalidates the buildings insurance, or that "
        "is illegal or immoral. Noise audible outside the Property between 11pm and "
        "7am is treated as a breach of this clause."])

    f += clause(4, "Access by the Landlord", [
        "The Landlord and anyone authorised by the Landlord may enter the Property at "
        "reasonable hours on at least twenty-four hours' written notice to inspect its "
        "condition, to carry out repairs, or to show it to a prospective tenant or "
        "purchaser during the final two months of the Term. In a genuine emergency, "
        "including fire, flood, or a suspected gas leak, entry may be made without "
        "notice and the Landlord shall tell the Tenants what happened as soon as "
        "practicable afterwards."])

    f += clause(5, "Assignment and subletting", [
        "The Tenants shall not assign, underlet, charge, or part with or share "
        "possession of the whole or any part of the Property. Taking in a lodger, "
        "listing the Property on a short-let platform, or allowing a person not named "
        "in Schedule 1 to reside at the Property for more than fourteen consecutive "
        "nights each amount to a breach of this clause."])

    f += clause(6, "Pets", [
        "The Tenants shall not keep any animal, bird, reptile or insect at the "
        "Property without the Landlord's prior written consent. Consent given for one "
        "animal does not extend to another. Where consent is given, the Tenants shall "
        "have the carpets and soft furnishings professionally cleaned at the end of "
        "the Term and shall produce the receipt."])

    f += clause(7, "Insurance", [
        "The Landlord shall insure the building and the Landlord's fixtures and "
        "fittings. That policy does not cover the Tenants' own possessions, and the "
        "Tenants are strongly advised to take out their own contents insurance. The "
        "Tenants shall not keep anything at the Property that would increase the "
        "premium or make the policy void, and shall pay any increase in premium "
        "caused by something they do."])

    f += clause(8, "Alterations and installations", [
        "The Tenants shall not alter the Property, install a satellite dish or an "
        "external aerial, change any lock, or add a smart doorbell or camera that "
        "records a shared hallway, without the Landlord's written consent. Where "
        "consent is given, the Tenants shall reinstate the Property at the end of the "
        "Term if the Landlord asks them to."])

    f += clause(9, "Garden and communal areas", [
        "The Tenants shall keep the allocated parking bay clear of anything other than "
        "a roadworthy and taxed vehicle, and shall not store a bicycle, pushchair or "
        "any other item in the communal stairwell, which is a protected fire escape "
        "route. Items left in the stairwell may be removed without notice and the cost "
        "of removal charged to the Tenants."])

    f += clause(10, "Deposit and deductions", [
        "At the end of the Term the Landlord may propose deductions from the Deposit "
        "for unpaid rent, for damage beyond fair wear and tear, for cleaning to bring "
        "the Property back to the standard recorded in the Inventory, and for the "
        "reasonable cost of removing anything the Tenants leave behind. The Landlord "
        "shall put any proposed deduction to the Tenants in writing with supporting "
        "evidence within ten working days of the end of the Term, and any part of the "
        "Deposit not in dispute shall be returned within the same period."],
        subs=[
            "A deduction may not be made for an item that the Inventory records as "
            "already damaged or missing at the start of the Term.",
            "Where the Tenants dispute a proposed deduction, either party may refer "
            "the dispute to the free adjudication service operated by the deposit "
            "scheme, and both parties agree to be bound by its decision.",
        ])

    f += clause(11, "Ending the tenancy", [
        "Either party may end the periodic tenancy that follows the Term by giving "
        "the other written notice expiring on the last day of a rental period, being "
        "at least one month's notice from the Tenants and at least two months' notice "
        "from the Landlord. Notice given by the Tenants must be given by all of them.",
        "On the last day of the Term or of any periodic tenancy the Tenants shall "
        "give up the Property with vacant possession, return every key, fob and "
        "remote control issued to them, and provide a forwarding address."])

    f += clause(12, "Notices", [
        "A notice under this agreement is validly given if it is delivered by hand, "
        "sent by first-class post, or sent to the email address each party has given "
        "in Schedule 1. A posted notice is treated as received on the second working "
        "day after posting, and an emailed notice on the next working day after it is "
        "sent, unless the sender receives a delivery failure message."])

    f.append(Spacer(1, 14))
    sig = [
        [Paragraph("<b>Signed by the Landlord's agent</b>", S["cell"]),
         Paragraph("<b>Signed by the Tenants</b>", S["cell"])],
        [Paragraph("<br/><br/>_______________________________<br/>"
                   "Aldergate Property Care Ltd<br/>Date: ______________", S["cell"]),
         Paragraph("<br/><br/>_______________________________<br/>"
                   "Dev Raghunathan<br/><br/>_______________________________<br/>"
                   "Priya Raghunathan<br/>Date: ______________", S["cell"])],
    ]
    st = Table(sig, colWidths=[85 * mm, 85 * mm])
    st.setStyle(TableStyle([
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ("LINEABOVE", (0, 0), (-1, 0), 0.5, RULE),
        ("TOPPADDING", (0, 0), (-1, -1), 6),
    ]))
    f.append(KeepTogether(st))

    doc.build(f)
    print("PDF  ->", path.name)


# --------------------------------------------------------------------------
# XLSX: multi-sheet operating budget
# --------------------------------------------------------------------------
from openpyxl import Workbook
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from openpyxl.utils import get_column_letter

HDR_FILL = PatternFill("solid", fgColor="1F3A5F")
SUB_FILL = PatternFill("solid", fgColor="DCE5F0")
TOT_FILL = PatternFill("solid", fgColor="F2E6CC")
HDR_FONT = Font(color="FFFFFF", bold=True, size=11)
THIN = Side(style="thin", color="B7C3D2")
BORDER = Border(left=THIN, right=THIN, top=THIN, bottom=THIN)
MONEY = '#,##0'


def style_header(ws, row, ncols):
    for c in range(1, ncols + 1):
        cell = ws.cell(row=row, column=c)
        cell.fill = HDR_FILL
        cell.font = HDR_FONT
        cell.border = BORDER
        cell.alignment = Alignment(horizontal="center" if c > 1 else "left",
                                   vertical="center", wrap_text=True)
    ws.row_dimensions[row].height = 30


def budget_xlsx():
    path = OUT / "Q3 Operating Budget FY26.xlsx"
    wb = Workbook()

    # -- Summary ----------------------------------------------------------
    ws = wb.active
    ws.title = "Summary"
    ws["A1"] = "Q3 FY26 Operating Budget — Consolidated"
    ws["A1"].font = Font(bold=True, size=14, color="1F3A5F")
    ws["A2"] = "Harrowgate Field Services Ltd  ·  prepared 8 Aug 2026  ·  all figures GBP"
    ws["A2"].font = Font(size=9, italic=True, color="666666")

    cols = ["Cost centre", "Jul budget", "Jul actual", "Aug budget", "Aug actual",
            "Sep budget", "Q3 budget", "Q3 forecast", "Variance", "Var %"]
    ws.append([])
    ws.append(cols)
    style_header(ws, 4, len(cols))

    rows = [
        ("4100 · Salaries and wages", 184200, 186940, 184200, 183110, 189600),
        ("4110 · Employer NI and pension", 31300, 31780, 31300, 31120, 32230),
        ("4200 · Contractors and agency", 42800, 51260, 38500, 44980, 36000),
        ("4300 · Depot rent and rates", 27450, 27450, 27450, 27450, 27450),
        ("4310 · Utilities and waste", 6120, 7345, 6120, 6890, 6400),
        ("4400 · Fleet lease and fuel", 33900, 36215, 33900, 38470, 34600),
        ("4410 · Vehicle maintenance", 9800, 12640, 9800, 8215, 10200),
        ("4500 · Plant hire", 18600, 15980, 21400, 22745, 19800),
        ("4510 · Consumables and PPE", 11250, 12085, 11250, 11640, 11500),
        ("4600 · Travel and subsistence", 8400, 9930, 8400, 10420, 8400),
        ("4700 · Software and licences", 14750, 14750, 14750, 15980, 15200),
        ("4710 · IT hardware", 5200, 3480, 5200, 9120, 5200),
        ("4800 · Marketing and tenders", 12300, 10240, 15600, 14870, 18900),
        ("4810 · Training and certification", 7600, 5120, 7600, 9340, 11400),
        ("4900 · Insurance", 16400, 16400, 16400, 16400, 16400),
        ("4910 · Professional fees", 9100, 11250, 9100, 7480, 9100),
        ("4920 · Bank and card charges", 1450, 1512, 1450, 1489, 1450),
        ("4990 · Contingency", 15000, 0, 15000, 0, 15000),
    ]
    r = 5
    for name, jb, ja, ab, aa, sb in rows:
        ws.cell(row=r, column=1, value=name)
        ws.cell(row=r, column=2, value=jb)
        ws.cell(row=r, column=3, value=ja)
        ws.cell(row=r, column=4, value=ab)
        ws.cell(row=r, column=5, value=aa)
        ws.cell(row=r, column=6, value=sb)
        ws.cell(row=r, column=7, value=f"=B{r}+D{r}+F{r}")
        ws.cell(row=r, column=8, value=f"=C{r}+E{r}+F{r}")
        ws.cell(row=r, column=9, value=f"=G{r}-H{r}")
        ws.cell(row=r, column=10, value=f"=IF(G{r}=0,0,(G{r}-H{r})/G{r})")
        r += 1

    last = r - 1
    ws.cell(row=r, column=1, value="Total operating cost").font = Font(bold=True)
    for c in range(2, 10):
        L = get_column_letter(c)
        ws.cell(row=r, column=c, value=f"=SUM({L}5:{L}{last})").font = Font(bold=True)
    ws.cell(row=r, column=10, value=f"=IF(G{r}=0,0,(G{r}-H{r})/G{r})").font = Font(bold=True)
    for c in range(1, 11):
        ws.cell(row=r, column=c).fill = TOT_FILL

    for rr in range(5, r + 1):
        for c in range(1, 11):
            cell = ws.cell(row=rr, column=c)
            cell.border = BORDER
            if c in (2, 3, 4, 5, 6, 7, 8, 9):
                cell.number_format = MONEY
            if c == 10:
                cell.number_format = '0.0%'
    for rr in range(5, last + 1, 2):
        for c in range(1, 11):
            if not ws.cell(row=rr, column=c).fill.fgColor.rgb == "00F2E6CC":
                ws.cell(row=rr, column=c).fill = PatternFill("solid", fgColor="F7F9FC")

    ws.column_dimensions["A"].width = 30
    for c in range(2, 11):
        ws.column_dimensions[get_column_letter(c)].width = 12.5
    ws.freeze_panes = "B5"

    # -- Depots -----------------------------------------------------------
    ws2 = wb.create_sheet("Depots")
    ws2["A1"] = "Cost by depot — Q3 FY26"
    ws2["A1"].font = Font(bold=True, size=13, color="1F3A5F")
    hdr = ["Depot", "Region", "Head-\ncount", "Jul", "Aug", "Sep (f)",
           "Q3 total", "Per head", "vs Q2", "Status"]
    ws2.append([]); ws2.append(hdr)
    style_header(ws2, 3, len(hdr))
    depots = [
        ("Willowmere Yard", "North", 34, 96420, 99180, 94300, "On plan"),
        ("Kestrel Bay Depot", "Coastal", 28, 81340, 88705, 83900, "Watch"),
        ("Thornhill Row", "North", 19, 54860, 52190, 55400, "On plan"),
        ("Marrow Lane Works", "Central", 41, 118900, 121460, 116750, "Watch"),
        ("Fenwick Sidings", "Central", 12, 33710, 31085, 34200, "On plan"),
        ("Alderbrook Store", "South", 9, 21480, 22940, 21900, "On plan"),
        ("Greywater Pumping", "Coastal", 15, 44290, 47630, 45100, "Over"),
        ("Saltmarsh Access", "Coastal", 7, 18960, 17240, 19400, "On plan"),
        ("Hollowfield Compound", "South", 23, 66180, 64920, 67300, "On plan"),
        ("Brackenridge Sub-depot", "North", 11, 29740, 33810, 30200, "Watch"),
    ]
    r = 4
    for name, region, hc, jul, aug, sep, status in depots:
        ws2.cell(row=r, column=1, value=name)
        ws2.cell(row=r, column=2, value=region)
        ws2.cell(row=r, column=3, value=hc)
        ws2.cell(row=r, column=4, value=jul)
        ws2.cell(row=r, column=5, value=aug)
        ws2.cell(row=r, column=6, value=sep)
        ws2.cell(row=r, column=7, value=f"=SUM(D{r}:F{r})")
        ws2.cell(row=r, column=8, value=f"=ROUND(G{r}/C{r},0)")
        ws2.cell(row=r, column=9, value=round((jul + aug + sep) / 292000 - 1, 4))
        ws2.cell(row=r, column=10, value=status)
        r += 1
    tot = r
    ws2.cell(row=tot, column=1, value="All depots").font = Font(bold=True)
    for c in [3, 4, 5, 6, 7]:
        L = get_column_letter(c)
        ws2.cell(row=tot, column=c, value=f"=SUM({L}4:{L}{tot-1})").font = Font(bold=True)
    for c in range(1, 11):
        ws2.cell(row=tot, column=c).fill = TOT_FILL
    for rr in range(4, tot + 1):
        for c in range(1, 11):
            cell = ws2.cell(row=rr, column=c)
            cell.border = BORDER
            if c in (4, 5, 6, 7, 8):
                cell.number_format = MONEY
            if c == 9:
                cell.number_format = '+0.0%;-0.0%'
            if c == 3:
                cell.alignment = Alignment(horizontal="center")
    ws2.column_dimensions["A"].width = 24
    ws2.column_dimensions["B"].width = 11
    for c in range(3, 11):
        ws2.column_dimensions[get_column_letter(c)].width = 11.5

    # -- Headcount --------------------------------------------------------
    ws3 = wb.create_sheet("Headcount")
    ws3["A1"] = "Establishment and vacancies"
    ws3["A1"].font = Font(bold=True, size=13, color="1F3A5F")
    hdr = ["Role", "Grade", "Budgeted\nFTE", "Filled", "Vacant", "Avg salary",
           "On-cost %", "Fully loaded", "Notes"]
    ws3.append([]); ws3.append(hdr)
    style_header(ws3, 3, len(hdr))
    roles = [
        ("Field technician", "T3", 62, 57, 5, 31400, 0.212, "3 offers out"),
        ("Senior field technician", "T4", 24, 24, 0, 38900, 0.212, ""),
        ("Depot supervisor", "M2", 11, 10, 1, 44700, 0.223, "Willowmere open"),
        ("Plant operator", "T3", 18, 16, 2, 33200, 0.212, "Agency covering"),
        ("Scheduler", "A2", 9, 9, 0, 29800, 0.205, ""),
        ("Compliance officer", "P2", 5, 4, 1, 41600, 0.218, "Interviewing"),
        ("Fleet coordinator", "A3", 4, 4, 0, 32500, 0.205, ""),
        ("Stores keeper", "A2", 8, 7, 1, 27900, 0.205, ""),
        ("HSE advisor", "P2", 3, 3, 0, 45300, 0.218, ""),
        ("Contract manager", "M3", 6, 5, 1, 56800, 0.231, "Backfill Q4"),
        ("Data analyst", "P1", 3, 2, 1, 38200, 0.212, "New for Q3"),
        ("Apprentice", "T1", 12, 11, 1, 19600, 0.198, "Sep intake"),
    ]
    r = 4
    for role, grade, fte, filled, vac, sal, onc, note in roles:
        ws3.cell(row=r, column=1, value=role)
        ws3.cell(row=r, column=2, value=grade)
        ws3.cell(row=r, column=3, value=fte)
        ws3.cell(row=r, column=4, value=filled)
        ws3.cell(row=r, column=5, value=vac)
        ws3.cell(row=r, column=6, value=sal)
        ws3.cell(row=r, column=7, value=onc)
        ws3.cell(row=r, column=8, value=f"=ROUND(D{r}*F{r}*(1+G{r}),0)")
        ws3.cell(row=r, column=9, value=note)
        r += 1
    for rr in range(4, r):
        for c in range(1, 10):
            cell = ws3.cell(row=rr, column=c)
            cell.border = BORDER
            if c in (6, 8):
                cell.number_format = MONEY
            if c == 7:
                cell.number_format = '0.0%'
            if c in (2, 3, 4, 5):
                cell.alignment = Alignment(horizontal="center")
    ws3.column_dimensions["A"].width = 24
    ws3.column_dimensions["I"].width = 18
    for c in range(2, 9):
        ws3.column_dimensions[get_column_letter(c)].width = 11

    # -- Capex ------------------------------------------------------------
    ws4 = wb.create_sheet("Capex")
    ws4["A1"] = "Capital requests awaiting approval"
    ws4["A1"].font = Font(bold=True, size=13, color="1F3A5F")
    hdr = ["Ref", "Item", "Depot", "Requested", "Approved", "Payback\n(months)",
           "Decision", "Sponsor"]
    ws4.append([]); ws4.append(hdr)
    style_header(ws4, 3, len(hdr))
    capex = [
        ("CX-2611", "Replace 3.5t tipper × 2", "Willowmere Yard", 74800, 74800, 41, "Approved", "R. Adeyemi"),
        ("CX-2612", "Mobile welfare unit", "Saltmarsh Access", 21400, 0, 28, "Deferred", "L. Brennan"),
        ("CX-2613", "Depot roof re-covering", "Marrow Lane Works", 96500, 88000, 0, "Approved", "M. Okonkwo"),
        ("CX-2614", "Telematics retrofit, 46 vehicles", "All", 38640, 38640, 14, "Approved", "T. Vasquez"),
        ("CX-2615", "Tool store racking", "Thornhill Row", 12900, 12900, 0, "Approved", "L. Brennan"),
        ("CX-2616", "Compressor overhaul", "Fenwick Sidings", 17250, 0, 22, "On hold", "R. Adeyemi"),
        ("CX-2617", "Site cameras, 4 compounds", "Various", 29300, 24000, 19, "Part", "S. Nakamura"),
        ("CX-2618", "EV charge points × 6", "Kestrel Bay Depot", 54100, 0, 63, "Under review", "T. Vasquez"),
        ("CX-2619", "Survey drone and licence", "Willowmere Yard", 8750, 8750, 11, "Approved", "S. Nakamura"),
    ]
    r = 4
    for row in capex:
        for i, v in enumerate(row, start=1):
            ws4.cell(row=r, column=i, value=v)
        r += 1
    ws4.cell(row=r, column=2, value="Total").font = Font(bold=True)
    ws4.cell(row=r, column=4, value=f"=SUM(D4:D{r-1})").font = Font(bold=True)
    ws4.cell(row=r, column=5, value=f"=SUM(E4:E{r-1})").font = Font(bold=True)
    for c in range(1, 9):
        ws4.cell(row=r, column=c).fill = TOT_FILL
    for rr in range(4, r + 1):
        for c in range(1, 9):
            cell = ws4.cell(row=rr, column=c)
            cell.border = BORDER
            if c in (4, 5):
                cell.number_format = MONEY
            if c == 6:
                cell.alignment = Alignment(horizontal="center")
    ws4.column_dimensions["A"].width = 10
    ws4.column_dimensions["B"].width = 28
    ws4.column_dimensions["C"].width = 20
    for c in range(4, 9):
        ws4.column_dimensions[get_column_letter(c)].width = 12.5

    # -- Assumptions ------------------------------------------------------
    ws5 = wb.create_sheet("Assumptions")
    ws5["A1"] = "Basis of preparation"
    ws5["A1"].font = Font(bold=True, size=13, color="1F3A5F")
    ws5.append([]); ws5.append(["Driver", "Value", "Source", "Last reviewed"])
    style_header(ws5, 3, 4)
    assum = [
        ("Pay award, effective 1 Jul", "3.20%", "Board minute 26/07", "12 Jun 2026"),
        ("Employer pension contribution", "6.00%", "Scheme rules v4", "01 Apr 2026"),
        ("Employer NI rate", "13.80%", "Payroll bureau", "06 Apr 2026"),
        ("Diesel, pence per litre", "148.5", "Fuel card 3-mo avg", "31 Jul 2026"),
        ("Fleet utilisation", "78%", "Telematics export", "31 Jul 2026"),
        ("Chargeable hours per FTE", "1,512", "Ops model v9", "20 May 2026"),
        ("Agency uplift on base rate", "21%", "Framework schedule", "01 Jan 2026"),
        ("Inflation applied to consumables", "2.40%", "Procurement index", "30 Jun 2026"),
        ("Bad debt provision", "0.90%", "Ledger 12-mo actual", "30 Jun 2026"),
        ("Contingency, % of opex", "3.00%", "Finance policy FP-11", "01 Apr 2026"),
    ]
    r = 4
    for row in assum:
        for i, v in enumerate(row, start=1):
            cell = ws5.cell(row=r, column=i, value=v)
            cell.border = BORDER
        r += 1
    ws5.column_dimensions["A"].width = 32
    ws5.column_dimensions["B"].width = 13
    ws5.column_dimensions["C"].width = 22
    ws5.column_dimensions["D"].width = 15

    wb.save(path)
    print("XLSX ->", path.name)


# --------------------------------------------------------------------------
# DOCX: formatted report
# --------------------------------------------------------------------------
import docx
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.shared import Pt, RGBColor, Inches


def report_docx():
    path = OUT / "Field Survey Report - Willowmere.docx"
    d = docx.Document()

    normal = d.styles["Normal"]
    normal.font.name = "Calibri"
    normal.font.size = Pt(11)
    normal.paragraph_format.space_after = Pt(7)
    normal.paragraph_format.line_spacing = 1.15

    for name, size, color in (("Heading 1", 18, "1F3A5F"), ("Heading 2", 14, "1F3A5F"),
                              ("Heading 3", 12, "2E5A88")):
        st = d.styles[name]
        st.font.name = "Calibri"
        st.font.size = Pt(size)
        st.font.bold = True
        st.font.color.rgb = RGBColor.from_string(color)

    t = d.add_paragraph()
    t.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = t.add_run("Willowmere Wetland Restoration")
    run.bold = True
    run.font.size = Pt(24)
    run.font.color.rgb = RGBColor.from_string("1F3A5F")

    s = d.add_paragraph()
    s.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r2 = s.add_run("Phase 2 Field Survey Report  ·  Contract WM-2026-114")
    r2.font.size = Pt(12)
    r2.font.color.rgb = RGBColor.from_string("666666")

    s2 = d.add_paragraph()
    s2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r3 = s2.add_run("Harrowgate Field Services Ltd  —  issued 6 August 2026  —  Revision C")
    r3.font.size = Pt(10)
    r3.italic = True
    r3.font.color.rgb = RGBColor.from_string("888888")

    d.add_heading("1. Summary of findings", level=1)
    d.add_paragraph(
        "Phase 2 of the Willowmere restoration was surveyed over eleven working days "
        "between 14 and 29 July 2026. The survey covered 46.3 hectares of former "
        "grazing marsh east of Marrow Lane, together with the 1.9 km of reprofiled "
        "ditch that now carries water from the Brackenridge outfall into the new "
        "central lagoon. Conditions were dry throughout, which gave good access to "
        "the southern compartments but depressed the counts for amphibians and for "
        "several of the wetland plant indicators.")
    d.add_paragraph(
        "The headline result is that the scheme is meeting its hydrological targets "
        "ahead of programme. Water levels in compartments 3 and 4 held within the "
        "design envelope for the whole of the survey window despite an unusually dry "
        "July, and the reedbed transplanted in March has taken across roughly "
        "three-quarters of its planted area. Two matters need attention before the "
        "winter: the bank slumping at chainage 0+840, and the persistent silt load "
        "arriving from the northern culvert.")

    d.add_heading("Key points", level=2)
    for b in [
        "Reed establishment across compartments 1 to 4 is 74% against a Phase 2 target of 60%.",
        "Standing water was retained in all four compartments through a 19-day dry spell.",
        "Bank slumping over a 22 m length at chainage 0+840 needs regrading before October.",
        "Silt arriving from the northern culvert is running at roughly 2.4 times the modelled rate.",
        "Breeding waders were recorded in compartment 2 for the first time since works began.",
        "No non-native invasive species were found within the survey boundary.",
    ]:
        d.add_paragraph(b, style="List Bullet")

    d.add_heading("2. Method", level=1)
    d.add_paragraph(
        "Vegetation was recorded using 2 m × 2 m quadrats on a stratified random "
        "grid, with twelve quadrats in each compartment and a further eight along the "
        "ditch margin. Percentage cover was estimated by eye against a printed "
        "reference card and cross-checked by a second surveyor on one quadrat in "
        "five. Water levels were read from the six dipwells installed in Phase 1 and "
        "from the two new wells at the lagoon inlet and outlet.")
    d.add_paragraph(
        "Bird records follow the standard territory-mapping approach over four visits, "
        "two of them starting before 06:00. Invertebrate sampling used sweep netting "
        "on the drier margins and a pond net in open water, with specimens identified "
        "in the field and released. Photographs were taken from the nine fixed points "
        "established at the outset so that later surveys remain comparable.")

    d.add_heading("3. Vegetation results", level=1)
    d.add_paragraph(
        "Cover figures for the four compartments are set out below. The figure for "
        "compartment 4 is depressed by the bare ground left after the March "
        "regrading, and is expected to close the gap on the others during the next "
        "growing season.")

    table = d.add_table(rows=1, cols=6)
    table.style = "Light Grid Accent 1"
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    hdr = table.rows[0].cells
    for i, h in enumerate(["Compartment", "Area (ha)", "Reed cover", "Open water",
                           "Bare ground", "Quadrats"]):
        hdr[i].text = h
        for p in hdr[i].paragraphs:
            for r in p.runs:
                r.bold = True
                r.font.size = Pt(10)
    data = [
        ("1 — North marsh", "11.4", "81%", "12%", "4%", "12"),
        ("2 — Central lagoon", "14.8", "76%", "19%", "3%", "12"),
        ("3 — South scrape", "9.7", "79%", "14%", "5%", "12"),
        ("4 — East regrade", "10.4", "58%", "9%", "28%", "12"),
        ("Ditch margin", "1.9", "63%", "22%", "11%", "8"),
    ]
    for row in data:
        cells = table.add_row().cells
        for i, v in enumerate(row):
            cells[i].text = v
            for p in cells[i].paragraphs:
                for r in p.runs:
                    r.font.size = Pt(10)

    d.add_paragraph()
    d.add_heading("4. Hydrology", level=1)
    d.add_paragraph(
        "The design intent was to hold water between 0.35 m and 0.70 m above datum in "
        "compartments 1 to 3 through the summer, drawing down gradually into September. "
        "Readings taken twice weekly show the range held between 0.41 m and 0.66 m "
        "across the survey window, which is comfortably inside the envelope and "
        "notably better than the equivalent period last year, when compartment 3 fell "
        "to 0.19 m and lost its standing water for nine days.")
    d.add_paragraph(
        "The one concern is sediment. Turbidity at the northern culvert averaged 118 "
        "NTU against a modelled 49 NTU, and the settlement pocket immediately "
        "downstream has already lost about a third of its designed capacity. Unless "
        "the pocket is cleared before the winter flows arrive, silt will begin to "
        "carry through into the lagoon and smother the reed rhizomes planted along "
        "the northern edge.")

    d.add_heading("5. Recommendations", level=1)
    for i, b in enumerate([
        "Regrade the slumped bank at chainage 0+840 to a 1:4 profile and reseed before the end of September.",
        "Clear the settlement pocket at the northern culvert and increase its capacity by roughly 40%.",
        "Retain the temporary stock fence around compartment 4 for one further season.",
        "Repeat the vegetation survey in late June 2027 so results are not skewed by another dry July.",
        "Add two dipwells at the eastern boundary, where the current network leaves a 400 m gap.",
    ], start=1):
        d.add_paragraph(b, style="List Number")

    d.add_heading("6. Limitations", level=1)
    d.add_paragraph(
        "This report covers the Phase 2 boundary only and does not assess the "
        "condition of the Phase 1 works upstream. The dry conditions during the "
        "survey window mean the amphibian and aquatic invertebrate counts should be "
        "treated as a minimum rather than a representative figure. Access to the "
        "north-eastern corner was restricted for two days by contractor works, and "
        "three quadrats there were recorded on the following visit instead.")

    d.save(path)
    print("DOCX ->", path.name)


# --------------------------------------------------------------------------
# PPTX: deck
# --------------------------------------------------------------------------
from pptx import Presentation
from pptx.dml.color import RGBColor as PRGB
from pptx.util import Emu, Inches as PIn, Pt as PPt
from pptx.enum.text import PP_ALIGN


def deck_pptx():
    path = OUT / "Willowmere Kickoff.pptx"
    prs = Presentation()
    prs.slide_width = PIn(13.333)
    prs.slide_height = PIn(7.5)
    blank = prs.slide_layouts[6]

    NAVY = PRGB(0x1F, 0x3A, 0x5F)
    TEAL = PRGB(0x2E, 0x8B, 0x84)
    SAND = PRGB(0xF6, 0xF1, 0xE6)
    GREY = PRGB(0x5A, 0x5A, 0x5A)

    def textbox(slide, l, t, w, h, text, size, bold=False, color=NAVY,
                align=PP_ALIGN.LEFT, font="Calibri"):
        tb = slide.shapes.add_textbox(PIn(l), PIn(t), PIn(w), PIn(h))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = text
        r.font.size = PPt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font
        return tb

    def band(slide, l, t, w, h, color):
        from pptx.enum.shapes import MSO_SHAPE
        sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PIn(l), PIn(t), PIn(w), PIn(h))
        sh.fill.solid()
        sh.fill.fore_color.rgb = color
        sh.line.fill.background()
        sh.shadow.inherit = False
        return sh

    # --- Slide 1: title
    s = prs.slides.add_slide(blank)
    band(s, 0, 0, 13.333, 7.5, SAND)
    band(s, 0, 0, 13.333, 1.55, NAVY)
    textbox(s, 0.85, 0.42, 11.5, 0.8, "HARROWGATE FIELD SERVICES", 15, True,
            PRGB(0xC8, 0xD8, 0xEA))
    textbox(s, 0.85, 2.15, 11.5, 1.4, "Willowmere Phase 3", 48, True, NAVY)
    textbox(s, 0.85, 3.25, 11.5, 0.9, "Project kickoff  ·  contract WM-2026-114",
            22, False, TEAL)
    band(s, 0.9, 4.25, 3.2, 0.06, TEAL)
    textbox(s, 0.85, 4.65, 11.5, 1.6,
            "46 hectares of restored grazing marsh, 1.9 km of reprofiled ditch, "
            "and a delivery window that closes on 31 March 2027.", 18, False, GREY)
    textbox(s, 0.85, 6.5, 11.5, 0.5,
            "6 August 2026   ·   Prepared for the Willowmere steering group",
            13, False, PRGB(0x8A, 0x8A, 0x8A))

    # --- Slide 2: agenda / content with bullets
    s = prs.slides.add_slide(blank)
    band(s, 0, 0, 13.333, 7.5, PRGB(0xFF, 0xFF, 0xFF))
    band(s, 0, 0, 0.28, 7.5, TEAL)
    textbox(s, 0.9, 0.55, 11.5, 0.9, "Where Phase 2 leaves us", 36, True, NAVY)
    band(s, 0.95, 1.5, 2.6, 0.05, TEAL)

    tb = s.shapes.add_textbox(PIn(0.9), PIn(1.9), PIn(6.4), PIn(4.9))
    tf = tb.text_frame
    tf.word_wrap = True
    points = [
        "Reed establishment reached 74%, against a 60% target",
        "Standing water held through a 19-day dry spell",
        "Breeding waders returned to compartment 2",
        "Silt from the northern culvert is 2.4× the modelled rate",
        "Bank slumping at chainage 0+840 needs regrading",
        "No invasive species found inside the boundary",
    ]
    for i, p_text in enumerate(points):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = "•  " + p_text
        r.font.size = PPt(19)
        r.font.color.rgb = PRGB(0x33, 0x33, 0x33)
        r.font.name = "Calibri"
        p.space_after = PPt(17)

    # stat cards
    cards = [("74%", "reed cover"), ("46.3 ha", "surveyed"), ("11", "survey days"),
             ("2.4×", "silt load")]
    for i, (big, small) in enumerate(cards):
        top = 1.95 + i * 1.22
        band(s, 7.9, top, 4.5, 1.05, PRGB(0xF1, 0xF5, 0xF9))
        band(s, 7.9, top, 0.09, 1.05, TEAL)
        textbox(s, 8.2, top + 0.08, 2.0, 0.6, big, 26, True, NAVY)
        textbox(s, 8.2, top + 0.58, 3.8, 0.4, small, 13, False, GREY)

    # --- Slide 3: milestones table-ish
    s = prs.slides.add_slide(blank)
    band(s, 0, 0, 13.333, 7.5, PRGB(0xFF, 0xFF, 0xFF))
    textbox(s, 0.9, 0.55, 11.5, 0.9, "Phase 3 milestones", 36, True, NAVY)
    band(s, 0.95, 1.5, 2.6, 0.05, TEAL)
    rows = [
        ("Sep 2026", "Bank regrade and reseed complete", "Site team"),
        ("Oct 2026", "Settlement pocket enlarged, silt trap in", "Contractor"),
        ("Dec 2026", "Winter water-level regime signed off", "Hydrology"),
        ("Feb 2027", "Eastern dipwell network extended", "Survey"),
        ("Mar 2027", "Phase 3 handover and final report", "All"),
    ]
    band(s, 0.9, 1.95, 11.5, 0.6, NAVY)
    textbox(s, 1.1, 2.05, 2.2, 0.4, "WHEN", 14, True, PRGB(0xFF, 0xFF, 0xFF))
    textbox(s, 3.5, 2.05, 6.5, 0.4, "MILESTONE", 14, True, PRGB(0xFF, 0xFF, 0xFF))
    textbox(s, 10.3, 2.05, 2.0, 0.4, "OWNER", 14, True, PRGB(0xFF, 0xFF, 0xFF))
    for i, (when, what, who) in enumerate(rows):
        top = 2.55 + i * 0.78
        if i % 2 == 0:
            band(s, 0.9, top, 11.5, 0.78, PRGB(0xF7, 0xF9, 0xFC))
        textbox(s, 1.1, top + 0.17, 2.2, 0.5, when, 16, True, TEAL)
        textbox(s, 3.5, top + 0.17, 6.6, 0.5, what, 16, False, PRGB(0x33, 0x33, 0x33))
        textbox(s, 10.3, top + 0.17, 2.2, 0.5, who, 16, False, GREY)

    prs.save(path)
    print("PPTX ->", path.name)


if __name__ == "__main__":
    lease_pdf()
    budget_xlsx()
    report_docx()
    deck_pptx()
    for p in sorted(OUT.iterdir()):
        print(f"   {p.stat().st_size/1024:8.1f} KB  {p.name}")
